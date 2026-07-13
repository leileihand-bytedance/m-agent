from __future__ import annotations

from pathlib import Path
from typing import Any, Iterable

from ..models import DocumentAsset, DocumentBlock, DocumentWarning


def parse_pptx_document(path: Path, *, asset_dir: Path) -> dict[str, object]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise RuntimeError("缺少 python-pptx，无法读取 PPTX 文件") from exc

    presentation = Presentation(str(path))
    asset_dir.mkdir(parents=True, exist_ok=True)
    blocks: list[DocumentBlock] = []
    assets: list[DocumentAsset] = []
    warnings: list[DocumentWarning] = []
    counter = 0

    def add_block(kind: str, text: str, location: str, shape: Any | None = None) -> None:
        nonlocal counter
        clean = text.strip()
        if not clean:
            return
        counter += 1
        blocks.append(
            DocumentBlock(
                block_id=f"block-{counter:05d}",
                kind=kind,
                text=clean,
                location=location,
                style=_shape_style(shape) if shape is not None else {},
                bbox=_shape_bbox(shape) if shape is not None else None,
            )
        )

    for slide_index, slide in enumerate(presentation.slides, 1):
        for shape_index, shape in enumerate(_walk_shapes(slide.shapes), 1):
            location = f"slide:{slide_index}/shape:{shape_index}"
            if getattr(shape, "has_table", False):
                rows = ["\t".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows]
                add_block("table", "\n".join(row for row in rows if row.strip()), location, shape)
                continue
            if getattr(shape, "has_chart", False):
                add_block("chart", _chart_text(shape.chart), location, shape)
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                asset = _save_picture(shape, slide_index=slide_index, shape_index=shape_index, asset_dir=asset_dir)
                if asset:
                    assets.append(asset)
                continue
            if getattr(shape, "has_text_frame", False):
                add_block("text", shape.text_frame.text, location, shape)

        if getattr(slide, "has_notes_slide", False):
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                add_block("notes", notes_frame.text, f"slide:{slide_index}/notes")

    metadata = {
        "title": presentation.core_properties.title or "",
        "subject": presentation.core_properties.subject or "",
        "author": presentation.core_properties.author or "",
        "slide_width": int(presentation.slide_width or 0),
        "slide_height": int(presentation.slide_height or 0),
    }
    return {
        "blocks": tuple(blocks),
        "assets": tuple(assets),
        "warnings": tuple(warnings),
        "page_count": len(presentation.slides),
        "metadata": metadata,
    }


def _walk_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        if hasattr(shape, "shapes"):
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


def _shape_bbox(shape: Any) -> tuple[int, int, int, int] | None:
    try:
        return (int(shape.left), int(shape.top), int(shape.width), int(shape.height))
    except (AttributeError, TypeError, ValueError):
        return None


def _shape_style(shape: Any) -> dict[str, object]:
    style: dict[str, object] = {"shape_name": str(getattr(shape, "name", "") or "")}
    frame = getattr(shape, "text_frame", None)
    if frame is not None:
        paragraph = next((item for item in frame.paragraphs if item.text.strip()), None)
        run = next((item for item in paragraph.runs if item.text.strip()), None) if paragraph else None
        if run is not None:
            font = run.font
            style.update(
                {
                    "font_name": str(font.name or ""),
                    "font_size_pt": float(font.size.pt) if font.size else None,
                    "bold": font.bold,
                    "italic": font.italic,
                }
            )
    return {key: value for key, value in style.items() if value not in (None, "")}


def _chart_text(chart: Any) -> str:
    parts: list[str] = []
    if getattr(chart, "has_title", False) and chart.chart_title.has_text_frame:
        parts.append(chart.chart_title.text_frame.text)
    for series in getattr(chart, "series", ()):
        name = str(getattr(series, "name", "") or "")
        values = []
        try:
            values = [str(value) for value in series.values]
        except Exception:
            pass
        text = "：".join(item for item in (name, "、".join(values)) if item)
        if text:
            parts.append(text)
    return "\n".join(parts)


def _save_picture(shape: Any, *, slide_index: int, shape_index: int, asset_dir: Path) -> DocumentAsset | None:
    try:
        image = shape.image
        extension = str(image.ext or "bin")
        filename = f"slide-{slide_index:03d}-shape-{shape_index:03d}.{extension}"
        target = asset_dir / filename
        target.write_bytes(image.blob)
        size = getattr(image, "size", (None, None))
        return DocumentAsset(
            asset_id=f"slide-{slide_index}-shape-{shape_index}",
            kind="image",
            location=f"slide:{slide_index}/shape:{shape_index}",
            path=str(target),
            content_type=str(getattr(image, "content_type", "") or ""),
            width=int(size[0]) if size and size[0] is not None else None,
            height=int(size[1]) if size and size[1] is not None else None,
        )
    except Exception:
        return None
